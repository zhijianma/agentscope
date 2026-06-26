# -*- coding: utf-8 -*-
"""PowerPoint (.pptx) file parser.

Walks the deck slide-by-slide and emits one :class:`Section` per
contiguous content block — adjacent text runs (and, by default,
tables) are merged into a single text section; embedded images are
emitted as their own :class:`DataBlock` sections.  Each section's
metadata carries the slide index (starting at 1) for later citation.

Mirrors the knob set of the v1 ``PowerPointReader``:
``include_image``, ``separate_table``, ``table_format``,
``slide_prefix``, ``slide_suffix``.  Chunking is **not** done here —
long text stays intact inside a section and is split downstream by a
:class:`~agentscope.rag.ChunkerBase`.
"""
import base64
import io
from typing import Any, Literal

from ..._logging import logger
from ...message import Base64Source, DataBlock, TextBlock
from .._document import Section
from ._base import ParserBase
from ._utils import (
    _guess_image_media_type,
    _table_to_json,
    _table_to_markdown,
)


def _extract_table_rows(table: Any) -> list[list[str]]:
    """Read a python-pptx table into a 2-D ``list[list[str]]``.

    Args:
        table (`Any`):
            The python-pptx ``Table`` object.

    Returns:
        `list[list[str]]`:
            One inner list per row; per-cell line breaks are
            normalised to ``\\n``.
    """
    rows: list[list[str]] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            text = cell.text.strip()
            text = text.replace("\r\n", "\n").replace("\r", "\n")
            cells.append(text)
        rows.append(cells)
    return rows


def _extract_image_bytes(shape: Any) -> bytes | None:
    """Return the embedded image bytes for a picture shape, or ``None``.

    Args:
        shape (`Any`):
            A python-pptx shape.  Non-picture shapes return ``None``.

    Returns:
        `bytes | None`:
            The raw image bytes, or ``None`` when ``shape`` is not a
            picture / the bytes are unreadable.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        picture_type = MSO_SHAPE_TYPE.PICTURE
    except ImportError:
        # MSO_SHAPE_TYPE.PICTURE numeric value used as the fallback
        # so the parser still works against pptx builds where the
        # enum import path has moved.
        picture_type = 13

    if shape.shape_type != picture_type:
        return None
    try:
        return shape.image.blob
    except Exception as e:  # pylint: disable=broad-except
        logger.warning("Failed to extract image from PPT shape: %s", e)
        return None


class PPTParser(ParserBase):
    """Parser for PowerPoint ``.pptx`` files.

    Slide order is preserved.  Within a slide, shapes are visited in
    document order and grouped into a minimum number of sections:

    - **Text** shapes and **table** shapes contribute to one running
      text section.  When ``separate_table=True`` a table closes the
      running section and starts a new one of its own.
    - **Picture** shapes emit a standalone :class:`Section` whose
      ``content`` is a :class:`DataBlock` holding the base64-encoded
      image bytes.

    ``slide_prefix`` / ``slide_suffix`` are wrapped around each
    slide's text content (the prefix is prepended to the slide's
    first text section, the suffix is appended to its last text
    section).  Use ``None`` on either to disable wrapping.
    """

    supported_media_types: list[str] = [
        "application/vnd.openxmlformats-officedocument.presentationml"
        ".presentation",
    ]

    @classmethod
    def supported_extensions(cls) -> list[str]:
        """Return ``[".pptx"]`` — the only format ``python-pptx``
        reads."""
        return [".pptx"]

    def __init__(
        self,
        include_image: bool = True,
        separate_table: bool = False,
        table_format: Literal["markdown", "json"] = "markdown",
        slide_prefix: str | None = "<slide index={index}>",
        slide_suffix: str | None = "</slide>",
    ) -> None:
        """Initialize the PowerPoint parser.

        Args:
            include_image (`bool`, defaults to ``True``):
                When ``True``, picture shapes are emitted as
                :class:`DataBlock` sections.  Set to ``False`` to keep
                a text-only index.
            separate_table (`bool`, defaults to ``False``):
                When ``True``, each table becomes its own text
                section, never merged with surrounding text.
            table_format (`Literal["markdown", "json"]`, defaults to
                ``"markdown"``):
                How to render tables.  ``"markdown"`` uses pipe-table
                syntax; ``"json"`` emits a JSON array prefixed with a
                ``<system-info>`` marker — choose JSON when cells
                contain newlines that would corrupt Markdown layout.
            slide_prefix (`str | None`, defaults to
                ``"<slide index={index}>"``):
                Prepended to the first text section of each slide.
                Supports the ``{index}`` placeholder (starting at 1).  Use
                ``None`` to disable.
            slide_suffix (`str | None`, defaults to ``"</slide>"``):
                Appended to the last text section of each slide.  Use
                ``None`` to disable.

        Raises:
            `ValueError`: If ``table_format`` is not one of
                ``"markdown"`` / ``"json"``.
        """
        if table_format not in ("markdown", "json"):
            raise ValueError(
                "The table_format must be one of 'markdown' or 'json', "
                f"got {table_format!r}.",
            )
        self.include_image = include_image
        self.separate_table = separate_table
        self.table_format = table_format
        self.slide_prefix = slide_prefix
        self.slide_suffix = slide_suffix

    async def parse(
        self,
        file: bytes | str,
        filename: str,
    ) -> list[Section]:
        """Parse a PPTX file into a list of :class:`Section` objects.

        Args:
            file (`bytes | str`):
                Either the raw PPTX bytes, or a filesystem path to
                the PPTX file.
            filename (`str`):
                The source filename, copied into each Section's
                :attr:`Section.source`.

        Returns:
            `list[Section]`:
                Sections in deck order.  Text sections carry
                ``{"slide": <starting at 1>}``; image sections add
                ``{"media_type": "image/..."}``.

        Raises:
            `FileNotFoundError`: If ``file`` is a ``str`` pointing to
                a path that does not exist.
            `ImportError`: If :mod:`python-pptx` is not installed.
            `ValueError`: If the bytes cannot be parsed.
        """
        if isinstance(file, str):
            with open(file, "rb") as fp:
                file = fp.read()

        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError(
                "Please install python-pptx to use the PowerPoint "
                "parser. You can install it by "
                "`pip install python-pptx` (or "
                "`pip install agentscope[rag]`).",
            ) from e

        try:
            prs = Presentation(io.BytesIO(file))
        except Exception as e:  # pylint: disable=broad-except
            raise ValueError(
                f"Failed to parse {filename!r} as PPTX: {e}",
            ) from e

        sections: list[Section] = []
        for slide_idx, slide in enumerate(prs.slides):
            sections.extend(
                self._parse_slide(slide, slide_idx, filename),
            )
        return sections

    # ------------------------------------------------------------------
    # Slide-level parsing
    # ------------------------------------------------------------------

    def _parse_slide(
        self,
        slide: Any,
        slide_idx: int,
        filename: str,
    ) -> list[Section]:
        """Walk one slide and return its ordered sections."""
        slide_no = slide_idx + 1
        prefix = (
            self.slide_prefix.format(index=slide_no)
            if self.slide_prefix is not None
            else ""
        )

        slide_sections: list[Section] = []
        # ``text_buffer`` accumulates the running text section; it is
        # flushed whenever an image arrives, the slide ends, or
        # ``separate_table`` forces a break around a table shape.
        text_buffer: list[str] = []

        def flush_text() -> None:
            if not text_buffer:
                return
            slide_sections.append(
                Section(
                    content=TextBlock(text="\n".join(text_buffer)),
                    source=filename,
                    metadata={"slide": slide_no},
                ),
            )
            text_buffer.clear()

        # Slide prefix lives at the very top of the first text section.
        if prefix:
            text_buffer.append(prefix)

        for shape in slide.shapes:
            # 1. Pictures — flush running text, emit a DataBlock section.
            if self.include_image:
                image_bytes = _extract_image_bytes(shape)
                if image_bytes is not None:
                    flush_text()
                    media_type = _guess_image_media_type(image_bytes)
                    data = base64.b64encode(image_bytes).decode("utf-8")
                    slide_sections.append(
                        Section(
                            content=DataBlock(
                                source=Base64Source(
                                    media_type=media_type,
                                    data=data,
                                ),
                                name=filename,
                            ),
                            source=filename,
                            metadata={
                                "slide": slide_no,
                                "media_type": media_type,
                            },
                        ),
                    )
                    continue

            # 2. Tables — render to text; optionally flush around them.
            if getattr(shape, "has_table", False):
                try:
                    rows = _extract_table_rows(shape.table)
                except Exception as e:  # pylint: disable=broad-except
                    logger.warning(
                        "Failed to extract table from slide %d: %s",
                        slide_no,
                        e,
                    )
                    continue
                rendered = (
                    _table_to_markdown(rows)
                    if self.table_format == "markdown"
                    else _table_to_json(rows)
                )
                if not rendered:
                    continue
                if self.separate_table:
                    flush_text()
                    slide_sections.append(
                        Section(
                            content=TextBlock(text=rendered),
                            source=filename,
                            metadata={"slide": slide_no},
                        ),
                    )
                else:
                    text_buffer.append(rendered)
                continue

            # 3. Text frames.
            if getattr(shape, "has_text_frame", False):
                try:
                    parts = [
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                        if para.text.strip()
                    ]
                except Exception as e:  # pylint: disable=broad-except
                    logger.warning(
                        "Failed to extract text from shape in slide %d: %s",
                        slide_no,
                        e,
                    )
                    continue
                if parts:
                    text_buffer.append("\n".join(parts))

        # Suffix goes onto the last text section of this slide.  If the
        # slide ends on an image (text_buffer empty, no prior text
        # section in this slide), create a small text-only carrier so
        # the suffix is preserved.
        if self.slide_suffix is not None:
            text_buffer.append(self.slide_suffix)

        flush_text()
        return slide_sections
