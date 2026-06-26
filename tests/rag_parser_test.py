# -*- coding: utf-8 -*-
"""Unit tests for the file parsers in :mod:`agentscope.rag._parser`.

PDF / PPTX fixtures are produced in-memory via :mod:`reportlab` and
:mod:`python-pptx` so the tests have no on-disk dependencies and can
run anywhere ``agentscope[rag]`` is installed.
"""
import base64
import io
import os
from unittest.async_case import IsolatedAsyncioTestCase

from utils import AnyString

from agentscope.rag import (
    ImageParser,
    PDFParser,
    PPTParser,
    TextParser,
)


# Smallest valid PNG ever — a single transparent pixel.  Used wherever
# a test needs "some image bytes" without depending on Pillow.
_PNG_PIXEL: bytes = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNgAAIAAAUA"
    "AarVyFEAAAAASUVORK5CYII=",
)
_PNG_PIXEL_B64: str = base64.b64encode(_PNG_PIXEL).decode("utf-8")


def _make_pdf(pages: list[str]) -> bytes:
    """Render the supplied page strings to a PDF in memory."""
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    for page in pages:
        pdf.drawString(72, 720, page)
        pdf.showPage()
    pdf.save()
    return buffer.getvalue()


def _make_pptx_simple(slides: list[str]) -> bytes:
    """Build a PPTX in memory with one text shape per slide."""
    from pptx import Presentation

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title-only layout
    for text in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = text
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _make_pptx_rich() -> bytes:
    """Build a richer PPTX with text, a table, and an embedded image.

    Slide 1: a single title text shape.
    Slide 2: a title text shape, a 2x2 table, and a trailing text shape.
    Slide 3: an embedded PNG image only (blank title).
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]

    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.title.text = "Hello"

    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.title.text = "Header"
    table = slide2.shapes.add_table(
        rows=2,
        cols=2,
        left=Inches(1),
        top=Inches(2),
        width=Inches(4),
        height=Inches(1),
    ).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    trailing = slide2.shapes.add_textbox(
        Inches(1),
        Inches(4),
        Inches(4),
        Inches(1),
    )
    trailing.text_frame.text = "Footer"

    slide3 = prs.slides.add_slide(blank_layout)
    slide3.shapes.title.text = ""  # blank title
    slide3.shapes.add_picture(
        io.BytesIO(_PNG_PIXEL),
        Inches(1),
        Inches(1),
        Inches(1),
        Inches(1),
    )

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class TextParserTest(IsolatedAsyncioTestCase):
    """Behavioural coverage for :class:`TextParser`."""

    async def test_decode_bytes(self) -> None:
        """Bytes are decoded with the configured encoding."""
        parser = TextParser()
        sections = await parser.parse(b"hello", "x.txt")
        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "hello",
                        "id": AnyString(),
                    },
                    "source": "x.txt",
                    "metadata": {},
                },
            ],
        )

    async def test_pre_decoded_string_round_trips(self) -> None:
        """``str`` inputs skip the decode step when no such file exists."""
        parser = TextParser()
        sections = await parser.parse("preset", "x.md")
        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "preset",
                        "id": AnyString(),
                    },
                    "source": "x.md",
                    "metadata": {},
                },
            ],
        )

    async def test_string_input_treated_as_path(self) -> None:
        """A ``str`` that names an existing file is read as a path."""
        import tempfile

        with tempfile.NamedTemporaryFile(
            mode="w",
            suffix=".txt",
            delete=False,
        ) as f:
            f.write("from disk")
            path = f.name
        try:
            parser = TextParser()
            sections = await parser.parse(path, "x.txt")
            self.assertEqual(len(sections), 1)
            self.assertEqual(sections[0].content.text, "from disk")
        finally:
            os.unlink(path)

    async def test_bad_encoding_raises_value_error(self) -> None:
        """An undecodable byte sequence surfaces a clear error."""
        parser = TextParser(encoding="ascii")
        with self.assertRaises(ValueError):
            await parser.parse(b"\xff\xfe", "bad.txt")


class PDFParserTest(IsolatedAsyncioTestCase):
    """Behavioural coverage for :class:`PDFParser`."""

    async def test_one_section_per_page(self) -> None:
        """Each page in the PDF round-trips to its own Section."""
        pdf_bytes = _make_pdf(["First page text", "Second page text"])
        parser = PDFParser()
        sections = await parser.parse(pdf_bytes, "demo.pdf")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "First page text\n",
                        "id": AnyString(),
                    },
                    "source": "demo.pdf",
                    "metadata": {"page": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "Second page text\n",
                        "id": AnyString(),
                    },
                    "source": "demo.pdf",
                    "metadata": {"page": 2},
                },
            ],
        )

    async def test_string_input_treated_as_path(self) -> None:
        """``str`` is interpreted as a filesystem path to the PDF."""
        import tempfile

        pdf_bytes = _make_pdf(["Hello"])
        with tempfile.NamedTemporaryFile(
            suffix=".pdf",
            delete=False,
        ) as f:
            f.write(pdf_bytes)
            path = f.name
        try:
            parser = PDFParser()
            sections = await parser.parse(path, "demo.pdf")
            self.assertEqual(len(sections), 1)
        finally:
            os.unlink(path)

    async def test_missing_path_raises_file_not_found(self) -> None:
        """A ``str`` pointing at a missing path raises FileNotFoundError."""
        parser = PDFParser()
        with self.assertRaises(FileNotFoundError):
            await parser.parse("/no/such/file.pdf", "x.pdf")

    async def test_invalid_bytes_raise_value_error(self) -> None:
        """Garbage bytes surface as :class:`ValueError`."""
        parser = PDFParser()
        with self.assertRaises(ValueError):
            await parser.parse(b"not a pdf", "broken.pdf")

    async def test_supported_extensions(self) -> None:
        """``.pdf`` is the only extension exposed to the file picker."""
        self.assertEqual(PDFParser.supported_extensions(), [".pdf"])


class ImageParserTest(IsolatedAsyncioTestCase):
    """Behavioural coverage for :class:`ImageParser`."""

    async def test_wraps_bytes_in_data_block(self) -> None:
        """A single Section is emitted with the image as a DataBlock."""
        parser = ImageParser()
        sections = await parser.parse(_PNG_PIXEL, "pixel.png")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "data",
                        "id": AnyString(),
                        "source": {
                            "type": "base64",
                            "data": _PNG_PIXEL_B64,
                            "media_type": "image/png",
                        },
                        "name": "pixel.png",
                    },
                    "source": "pixel.png",
                    "metadata": {"media_type": "image/png"},
                },
            ],
        )

    async def test_media_type_sniffed_from_jpeg(self) -> None:
        """JPEG magic bytes yield ``image/jpeg``."""
        jpeg_bytes = b"\xff\xd8\xff\xe0rest-of-jpeg"
        parser = ImageParser()
        sections = await parser.parse(jpeg_bytes, "x.jpg")
        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "data",
                        "id": AnyString(),
                        "source": {
                            "type": "base64",
                            "data": base64.b64encode(jpeg_bytes).decode(
                                "utf-8",
                            ),
                            "media_type": "image/jpeg",
                        },
                        "name": "x.jpg",
                    },
                    "source": "x.jpg",
                    "metadata": {"media_type": "image/jpeg"},
                },
            ],
        )

    async def test_string_input_treated_as_path(self) -> None:
        """``str`` is interpreted as a filesystem path to the image."""
        import tempfile

        with tempfile.NamedTemporaryFile(
            suffix=".png",
            delete=False,
        ) as f:
            f.write(_PNG_PIXEL)
            path = f.name
        try:
            parser = ImageParser()
            sections = await parser.parse(path, "pixel.png")
            self.assertEqual(len(sections), 1)
        finally:
            os.unlink(path)

    async def test_missing_path_raises_file_not_found(self) -> None:
        """A ``str`` pointing at a missing path raises FileNotFoundError."""
        parser = ImageParser()
        with self.assertRaises(FileNotFoundError):
            await parser.parse("/no/such/file.png", "x.png")


class PPTParserTest(IsolatedAsyncioTestCase):
    """Behavioural coverage for :class:`PPTParser`."""

    async def test_simple_deck_text_only(self) -> None:
        """A simple text-only deck round-trips through wrapping tags."""
        pptx_bytes = _make_pptx_simple(["Alpha", "Beta"])
        parser = PPTParser(include_image=False)
        sections = await parser.parse(pptx_bytes, "demo.pptx")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=1>\nAlpha\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "demo.pptx",
                    "metadata": {"slide": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=2>\nBeta\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "demo.pptx",
                    "metadata": {"slide": 2},
                },
            ],
        )

    async def test_without_slide_tags(self) -> None:
        """Disabling prefix/suffix removes the slide tags entirely."""
        pptx_bytes = _make_pptx_simple(["Alpha"])
        parser = PPTParser(
            include_image=False,
            slide_prefix=None,
            slide_suffix=None,
        )
        sections = await parser.parse(pptx_bytes, "demo.pptx")
        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "Alpha",
                        "id": AnyString(),
                    },
                    "source": "demo.pptx",
                    "metadata": {"slide": 1},
                },
            ],
        )

    async def test_table_merges_with_surrounding_text_by_default(
        self,
    ) -> None:
        """``separate_table=False`` merges the table into the running
        text section."""
        pptx_bytes = _make_pptx_rich()
        parser = PPTParser(include_image=False, separate_table=False)
        sections = await parser.parse(pptx_bytes, "rich.pptx")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=1>\nHello\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": (
                            "<slide index=2>\n"
                            "Header\n"
                            "| A | B |\n"
                            "| --- | --- |\n"
                            "| 1 | 2 |\n\n"
                            "Footer\n"
                            "</slide>"
                        ),
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=3>\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 3},
                },
            ],
        )

    async def test_table_separated_when_separate_table_true(self) -> None:
        """``separate_table=True`` flushes the running text around the
        table."""
        pptx_bytes = _make_pptx_rich()
        parser = PPTParser(include_image=False, separate_table=True)
        sections = await parser.parse(pptx_bytes, "rich.pptx")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=1>\nHello\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=2>\nHeader",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": (
                            "| A | B |\n" "| --- | --- |\n" "| 1 | 2 |\n"
                        ),
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "Footer\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=3>\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 3},
                },
            ],
        )

    async def test_image_emits_data_block(self) -> None:
        """An embedded picture becomes its own DataBlock section."""
        pptx_bytes = _make_pptx_rich()
        parser = PPTParser(include_image=True, separate_table=True)
        sections = await parser.parse(pptx_bytes, "rich.pptx")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=1>\nHello\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=2>\nHeader",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": (
                            "| A | B |\n" "| --- | --- |\n" "| 1 | 2 |\n"
                        ),
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "Footer\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=3>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 3},
                },
                {
                    "content": {
                        "type": "data",
                        "id": AnyString(),
                        "source": {
                            "type": "base64",
                            "data": _PNG_PIXEL_B64,
                            "media_type": "image/png",
                        },
                        "name": "rich.pptx",
                    },
                    "source": "rich.pptx",
                    "metadata": {
                        "slide": 3,
                        "media_type": "image/png",
                    },
                },
                {
                    "content": {
                        "type": "text",
                        "text": "</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 3},
                },
            ],
        )

    async def test_table_json_format(self) -> None:
        """``table_format="json"`` emits the JSON marker payload."""
        pptx_bytes = _make_pptx_rich()
        parser = PPTParser(
            include_image=False,
            separate_table=True,
            table_format="json",
        )
        sections = await parser.parse(pptx_bytes, "rich.pptx")

        self.assertEqual(
            [s.model_dump() for s in sections],
            [
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=1>\nHello\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 1},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=2>\nHeader",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": (
                            "<system-info>A table loaded as a JSON "
                            "array:</system-info>\n"
                            '[["A", "B"], ["1", "2"]]'
                        ),
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "Footer\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 2},
                },
                {
                    "content": {
                        "type": "text",
                        "text": "<slide index=3>\n</slide>",
                        "id": AnyString(),
                    },
                    "source": "rich.pptx",
                    "metadata": {"slide": 3},
                },
            ],
        )

    async def test_table_format_validation(self) -> None:
        """Unknown ``table_format`` raises :class:`ValueError`."""
        with self.assertRaises(ValueError):
            PPTParser(table_format="csv")  # type: ignore[arg-type]

    async def test_string_input_treated_as_path(self) -> None:
        """``str`` is interpreted as a filesystem path to the PPTX."""
        import tempfile

        pptx_bytes = _make_pptx_simple(["Alpha"])
        with tempfile.NamedTemporaryFile(
            suffix=".pptx",
            delete=False,
        ) as f:
            f.write(pptx_bytes)
            path = f.name
        try:
            parser = PPTParser(include_image=False)
            sections = await parser.parse(path, "demo.pptx")
            self.assertEqual(len(sections), 1)
        finally:
            os.unlink(path)

    async def test_missing_path_raises_file_not_found(self) -> None:
        """A ``str`` pointing at a missing path raises FileNotFoundError."""
        parser = PPTParser()
        with self.assertRaises(FileNotFoundError):
            await parser.parse("/no/such/file.pptx", "x.pptx")
